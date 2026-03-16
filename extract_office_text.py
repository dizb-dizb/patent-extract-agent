"""
Extract text from Office files into plain .txt for further analysis.

Supported:
- .docx via python-docx
- .pptx via python-pptx (optional)

.doc (legacy Word) is not reliably supported without external converters.
"""

from __future__ import annotations

from pathlib import Path


def extract_docx(path: Path) -> str:
    from docx import Document  # type: ignore

    doc = Document(str(path))
    lines: list[str] = []
    for para in doc.paragraphs:
        t = (para.text or "").strip()
        if t:
            lines.append(t)
    return "\n".join(lines).strip()


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:  # pragma: no cover
        raise RuntimeError("python-pptx not installed") from e

    prs = Presentation(str(path))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if not text:
                continue
            for line in str(text).splitlines():
                line = line.strip()
                if line:
                    lines.append(line)
    return "\n".join(lines).strip()


def main() -> None:
    targets = [
        Path(r"e:\dizb\Documents\OneDrive\Desktop\2_任务书-基于小样本学习的专利术语体抽取算法设计与实现.docx"),
        Path(r"e:\dizb\Documents\OneDrive\Desktop\22281153-邓卓斌-开题报告.doc"),
        Path(r"e:\dizb\Documents\OneDrive\Desktop\22281153-邓卓斌-开题答辩.pptx"),
    ]

    out_dir = Path(__file__).resolve().parent / "_office_text"
    out_dir.mkdir(parents=True, exist_ok=True)

    for p in targets:
        if not p.exists():
            print(f"[skip] not found: {p}")
            continue
        suf = p.suffix.lower()
        out = out_dir / (p.stem + ".txt")
        try:
            if suf == ".docx":
                text = extract_docx(p)
            elif suf == ".pptx":
                text = extract_pptx(p)
            elif suf == ".doc":
                raise RuntimeError(".doc is legacy format; please export as .docx or .pdf")
            else:
                raise RuntimeError(f"unsupported suffix: {suf}")
            out.write_text(text, encoding="utf-8")
            print(f"[ok] {p.name} -> {out} ({len(text)} chars)")
        except Exception as e:
            print(f"[fail] {p.name}: {e}")


if __name__ == "__main__":
    main()

